import os
import sys

try:
    from pdf2docx import Converter
except ImportError as e:
    print(f"Warning: pdf2docx not available: {e}")
    Converter = None

try:
    from docx2pdf import convert
except ImportError as e:
    print(f"Warning: docx2pdf not available: {e}")
    convert = None

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print(f"Warning: python-pptx not available: {e}")
    Presentation = None

try:
    import PyPDF2
except ImportError as e:
    print(f"Warning: PyPDF2 not available: {e}")
    PyPDF2 = None

def pdf_to_word(pdf_file: str, word_file: str = None):
    """Convert PDF to Word DOCX."""
    if not Converter:
        raise ImportError("pdf2docx library is not installed")
    
    if not word_file:
        word_file = os.path.splitext(pdf_file)[0] + ".docx"
    
    try:
        cv = Converter(pdf_file)
        cv.convert(word_file, start=0, end=None)
        cv.close()
        return word_file
    except Exception as e:
        print(f"Error converting PDF to Word: {e}")
        raise

def word_to_pdf(word_file: str, pdf_file: str = None):
    """Convert Word DOCX to PDF (Windows only if using docx2pdf)."""
    if not pdf_file:
        pdf_file = os.path.splitext(word_file)[0] + ".pdf"
    
    if convert and sys.platform == 'win32':
        try:
            convert(word_file, pdf_file)
            return pdf_file
        except Exception as e:
            print(f"Error converting Word to PDF: {e}")
            raise
    else:
        raise NotImplementedError("Conversion from Word to PDF is only supported on Windows environments (requires Microsoft Word).")
    return pdf_file

def pdf_to_ppt(pdf_file: str, ppt_file: str = None):
    """Convert PDF text to PowerPoint slides."""
    if not Presentation or not PyPDF2:
        raise ImportError("python-pptx or PyPDF2 library is not installed")
    
    if not ppt_file:
        ppt_file = os.path.splitext(pdf_file)[0] + ".pptx"
        
    try:
        prs = Presentation()
        reader = PyPDF2.PdfReader(pdf_file)
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title to generic page number or first line
                title = slide.shapes.title
                title.text = f"Page {i + 1}"
                
                content = slide.placeholders[1]
                content.text = text
                
        prs.save(ppt_file)
        return ppt_file
    except Exception as e:
        print(f"Error converting PDF to PPT: {e}")
        return None


